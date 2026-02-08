"""PowerPoint (.pptx) document parser."""

from typing import List

from cat_agent.log import logger
from cat_agent.tools.parsers.base import clean_paragraph


def parse_ppt(path: str, extract_image: bool = False) -> List[dict]:
    if extract_image:
        raise ValueError('Currently, extracting images is not supported!')

    from pptx import Presentation
    from pptx.exc import PackageNotFoundError

    try:
        ppt = Presentation(path)
    except PackageNotFoundError as ex:
        logger.warning(ex)
        return []

    doc = []
    for slide_number, slide in enumerate(ppt.slides):
        page = {'page_num': slide_number + 1, 'content': []}

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ''.join(run.text for run in paragraph.runs)
                    paragraph_text = clean_paragraph(paragraph_text)
                    if paragraph_text.strip():
                        page['content'].append({'text': paragraph_text})

            if shape.has_table:
                tbl = []
                for row in shape.table.rows:
                    tbl.append('|' + '|'.join([cell.text for cell in row.cells]) + '|')
                tbl = '\n'.join(tbl)
                page['content'].append({'table': tbl})

        doc.append(page)
    return doc
